"""
Text extraction from workflow output files.

Supports .pdf (PyMuPDF), .docx (python-docx), .xlsx (openpyxl),
.pptx (python-pptx), .csv (csv stdlib), .html (existing html_to_text),
.zip (recursive). All imports are guarded so a missing library degrades
to '' rather than raising, consistent with the rest of the RAG subsystem.
"""
from __future__ import annotations

import base64
import csv
import io
import logging
import zipfile
from pathlib import Path

log = logging.getLogger('app.rag.file_extractor')


def extract_text_from_bytes(data: bytes,
                             mime_type: str = '',
                             file_ext: str = '') -> str:
    """Dispatch on file_ext / mime_type and return extracted plain text."""
    ext  = (file_ext or '').lower().lstrip('.')
    mime = (mime_type or '').lower()

    if ext == 'pdf' or 'pdf' in mime:
        return _pdf(data)
    if ext == 'docx' or ('openxmlformats' in mime and 'word' in mime):
        return _docx(data)
    if ext == 'xlsx' or ('openxmlformats' in mime and 'sheet' in mime):
        return _xlsx(data)
    if ext == 'pptx' or ('openxmlformats' in mime and 'presentation' in mime):
        return _pptx(data)
    if ext == 'csv' or mime == 'text/csv':
        return _csv(data)
    if ext in ('html', 'htm') or 'html' in mime:
        return _html(data)
    if ext == 'zip' or mime in ('application/zip', 'application/x-zip-compressed'):
        return _zip(data)
    if ext in ('txt', 'md') or mime.startswith('text/'):
        return data.decode('utf-8', errors='replace')
    try:
        return data.decode('utf-8', errors='replace')
    except Exception:
        return ''


def extract_from_base64(data_b64: str,
                        mime_type: str = '',
                        file_ext: str = '') -> str:
    """Convenience: decode base64 then delegate to extract_text_from_bytes."""
    if not data_b64:
        return ''
    try:
        raw = base64.b64decode(data_b64)
    except Exception as exc:
        log.warning('[EXTRACTOR] base64 decode failed: %s', exc)
        return ''
    return extract_text_from_bytes(raw, mime_type=mime_type, file_ext=file_ext)


# ── format handlers ──────────────────────────────────────────────────

def _pdf(data: bytes) -> str:
    try:
        import pymupdf                                      # PyMuPDF >= 1.24
        doc = pymupdf.open(stream=data, filetype='pdf')
        return '\n\n'.join(page.get_text() for page in doc).strip()
    except ImportError:
        pass
    try:
        import fitz                                         # older PyMuPDF
        doc = fitz.open(stream=data, filetype='pdf')
        return '\n\n'.join(page.get_text() for page in doc).strip()
    except ImportError:
        pass
    log.warning('[EXTRACTOR] PDF extraction requires PyMuPDF')
    return ''


def _docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n\n'.join(paras)
    except ImportError:
        log.warning('[EXTRACTOR] DOCX extraction requires python-docx')
        return ''


def _xlsx(data: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join('' if v is None else str(v) for v in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                parts.append(f'Sheet: {ws.title}\n' + '\n'.join(rows))
        return '\n\n'.join(parts)
    except ImportError:
        log.warning('[EXTRACTOR] XLSX extraction requires openpyxl')
        return ''


_XLSX_VIEW_MAX_ROWS = 500
_XLSX_VIEW_MAX_COLS = 50


def xlsx_to_html(data: bytes) -> str:
    """Render an .xlsx workbook as an HTML table for the document VIEWER
    (routes/agents.py's document/<id>/view) — deliberately server-side
    with openpyxl (already a vetted dependency of this same module,
    read_only + no formula evaluation) rather than a client-side JS xlsx
    parser: the well-known 'xlsx'/SheetJS npm package has unpatched high-
    severity prototype-pollution/ReDoS advisories, which is a real risk
    for a feature whose whole job is parsing files uploaded by (or
    forwarded from) a client. Capped at 500 rows / 50 cols per sheet —
    a preview, not a spreadsheet engine; use the original-file download
    for anything larger. Multiple sheets are rendered as separate tables
    with a heading each."""
    import html as _html_mod
    try:
        import openpyxl
    except ImportError:
        log.warning('[EXTRACTOR] XLSX view requires openpyxl')
        return '<p>Preview unavailable — openpyxl is not installed.</p>'
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        rows_html = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= _XLSX_VIEW_MAX_ROWS:
                rows_html.append('<tr><td colspan="99"><em>… truncated at '
                                 f'{_XLSX_VIEW_MAX_ROWS} rows — download the original for the rest</em></td></tr>')
                break
            cells = row[:_XLSX_VIEW_MAX_COLS]
            cells_html = ''.join(
                '<td>' + _html_mod.escape('' if v is None else str(v)) + '</td>' for v in cells
            )
            if cells_html.strip('<td></td>'):
                rows_html.append(f'<tr>{cells_html}</tr>')
        if rows_html:
            parts.append(f'<h4>{_html_mod.escape(ws.title)}</h4><table>{"".join(rows_html)}</table>')
    return '\n'.join(parts) or '<p>This spreadsheet has no readable cells.</p>'


def _pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            lines = []
            for shape in slide.shapes:
                if getattr(shape, 'has_text_frame', False) and shape.text_frame.text.strip():
                    lines.append(shape.text_frame.text.strip())
                if getattr(shape, 'has_table', False):
                    for row in shape.table.rows:
                        row_text = '\t'.join(c.text for c in row.cells)
                        if row_text.strip():
                            lines.append(row_text)
            if getattr(slide, 'has_notes_slide', False) and slide.notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or '').strip()
                if notes:
                    lines.append(f'Notes: {notes}')
            if lines:
                slides.append(f'Slide {i}\n' + '\n'.join(lines))
        return '\n\n'.join(slides)
    except ImportError:
        log.warning('[EXTRACTOR] PPTX extraction requires python-pptx')
        return ''


def _csv(data: bytes) -> str:
    try:
        text = data.decode('utf-8', errors='replace')
        reader = csv.reader(io.StringIO(text))
        return '\n'.join('\t'.join(row) for row in reader)
    except Exception as exc:
        log.warning('[EXTRACTOR] CSV extraction failed: %s', exc)
        return ''


def _html(data: bytes) -> str:
    from app.services.rag.chunking import html_to_text
    return html_to_text(data.decode('utf-8', errors='replace'))


def _zip(data: bytes) -> str:
    parts = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for name in zf.namelist():
                if name.endswith('/'):
                    continue
                ext = Path(name).suffix.lower().lstrip('.')
                try:
                    entry_bytes = zf.read(name)
                    text = extract_text_from_bytes(entry_bytes, file_ext=ext)
                    if text.strip():
                        parts.append(f'--- {name} ---\n{text}')
                except Exception as exc:
                    log.debug('[EXTRACTOR] ZIP entry %s skipped: %s', name, exc)
    except Exception as exc:
        log.warning('[EXTRACTOR] ZIP extraction failed: %s', exc)
    return '\n\n'.join(parts)
